"""
export/export_service.py
========================
Generates DOCX and PPTX exports from the last analysis result.
Called by POST /export in main.py.
"""
from __future__ import annotations

import io
from typing import Optional

from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt


def generate_docx(request) -> io.BytesIO:
    """Generate a Word document from the export request."""
    doc = Document()

    # Title
    title = doc.add_heading("RetailCo Campaign Intelligence Report", level=0)
    title.runs[0].font.color.rgb = RGBColor(0x1a, 0x56, 0xdb)

    # Summary
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(request.summary or "No summary available.")

    # Rich text
    if request.rich_text:
        doc.add_heading("Detailed Analysis", level=1)
        for line in request.rich_text.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            else:
                # Strip markdown bold
                line = line.replace("**", "")
                doc.add_paragraph(line)

    # Table
    if request.table_data and request.table_data.columns:
        doc.add_heading("Data Table", level=1)
        table = doc.add_table(
            rows=1 + len(request.table_data.rows),
            cols=len(request.table_data.columns),
        )
        table.style = "Table Grid"
        # Header
        for i, col in enumerate(request.table_data.columns):
            cell = table.rows[0].cells[i]
            cell.text = col
            cell.paragraphs[0].runs[0].bold = True
        # Rows
        for r_idx, row in enumerate(request.table_data.rows):
            for c_idx, val in enumerate(row):
                table.rows[r_idx + 1].cells[c_idx].text = str(val)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def generate_pptx(request) -> io.BytesIO:
    """Generate a PowerPoint presentation from the export request."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Campaign Intelligence Report"
    slide.placeholders[1].text = "RetailCo AI Analytics"

    # Slide 2 — Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = request.summary or "No summary available."
    tf.word_wrap = True

    # Slide 3 — Recommendations
    if request.chart_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = request.chart_title or "Campaign Performance"
        tf = slide.placeholders[1].text_frame
        tf.text = "Key Metrics:"
        for item in (request.chart_data or [])[:8]:
            p = tf.add_paragraph()
            p.text = f"• {item.get('name', '')}: {item.get('value', '')} ({item.get('metric', '')})"

    # Slide 4 — Table
    if request.table_data and request.table_data.columns:
        slide = prs.slides.add_slide(blank_layout)
        tf_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf_box.text_frame.text = "Data Summary"
        tf_box.text_frame.paragraphs[0].runs[0].font.size = PPTPt(24)
        tf_box.text_frame.paragraphs[0].runs[0].font.bold = True

        cols = len(request.table_data.columns)
        rows = min(len(request.table_data.rows) + 1, 10)
        tbl  = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0),
                                       Inches(9.0), Inches(5.0)).table
        for c, col in enumerate(request.table_data.columns):
            tbl.cell(0, c).text = col
        for r, row in enumerate(request.table_data.rows[:9]):
            for c, val in enumerate(row[:cols]):
                tbl.cell(r + 1, c).text = str(val)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
